#!/usr/bin/env python3
# -*- coding: utf-8 -*-
#
# USAGE: ./gen_ppt.py 2014.1R 2014.5R 2014.6R
#
# Install (Ubuntu 15.04) : python 3.4.3이 깔려 있는 상태에서 아래처럼
# $ sudo apt-get install libxml2-dev libxslt-dev libjpeg-dev \
#   libfreetype6 libfreetype6-dev zlib1g-dev
# $ pip3 install python-pptx
#

import sys
from getpass import getpass
import xmlrpc.client
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Pt

###################################
# 여기를 수정해 주세요
trac_url = '172.16.33.2:8080/trac'
name = '장세연'
team = 'I랩'
###################################


def set_text(cell, text):
    frame = cell.text_frame
    frame.paragraphs[0].font.size = Pt(9)
    frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    frame.paragraphs[0].text = text


def process_set_text(cell, texts, ticket_no, ticket, this_round):
    round = this_round[5:-1]
    site = ''
    project = ''
    prev = texts[-1] if texts else ''
    type = 'V' if ticket[3]['type'] == '휴가' else 'W'

    lines = ticket[3]['description'].splitlines(True)
    for l in lines:
        if '- 고객사: ' in l:
            site = l[7:-1]
        elif '- 프로젝트: ' in l:
            project = l[8:-1]
    if prev == '실행자:  ':
        set_text(cell, name)
    elif prev == '팀명 :':
        set_text(cell, team)
    elif prev == '고객사:':
        set_text(cell, site)
    elif prev == '프로젝트:':
        set_text(cell, project)
    elif prev == '업무내용:':
        set_text(cell, ticket[3]['summary'])
    elif prev == '2015년':
        set_text(cell, round)
    elif prev == '[#':
        set_text(cell, str(ticket_no))
    elif len(texts) > 10:
        if texts[-6] == type:
            set_text(cell, 'O')
        # 시작/예정/완료
        if texts[-1] == 'X' or texts[-4] == 'X' or texts[-6] == 'X':
            set_text(cell, 'O')


def main(trac_id, trac_pass, this_round):
    pptfile = '%s_tickets.pptx' % this_round
    print('\n' + pptfile)
    trac_rpc_url = 'http://%s:%s@%s/login/rpc' % (trac_id, trac_pass, trac_url)
    server = xmlrpc.client.ServerProxy(trac_rpc_url)
    tlist = server.ticket.query('owner=%s&due_date^=%s' % (name, this_round))
    tickets = []
    for t in tlist:
        ticket = server.ticket.get(t)
        print('ticket collected. #%d %s' % (t, ticket[3]['summary']))
        tickets.append(ticket)

    prs = Presentation('ticket.pptx')
    cell_texts = []
    t_idx = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for row in shape.table.rows:
                for cell in row.cells:
                    process_set_text(cell, cell_texts, tlist[t_idx],
                                     tickets[t_idx], this_round)
                    cell_texts.append(cell.text_frame.text)

                    # 한 티켓의 입력이 끝날때, 다음 티켓 입력할지 종료할지 결정
                    if cell_texts[-1] == '실행자':
                        if len(tlist) - 1 <= t_idx:
                            prs.save(pptfile)
                            return
                        else:
                            t_idx = t_idx + 1


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('USAGE: ./gen_ppt.py 2014.1R 2014.5R 2014.6R')
        sys.exit()

    trac_id = input('trac username:')
    trac_pass = getpass('trac password:')

    for arg in sys.argv:
        if arg == sys.argv[0]:
            continue
        main(trac_id, trac_pass, arg)
